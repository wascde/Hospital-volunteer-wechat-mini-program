from pptx import Presentation
prs = Presentation(r'D:\vollunteercloundapp\医院志愿者助手_项目汇报.pptx')
print(f'页数: {len(prs.slides)}')
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    texts.append(t[:50])
    first = texts[0] if texts else '(空)'
    print(f'  {i}. {first}')
