#!/usr/bin/env python3
"""
生成「医院志愿者助手」微信小程序项目汇报 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== 配色方案 ==========
BLUE = RGBColor(0x1A, 0x73, 0xE8)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    """填充纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK, spacing=Pt(8)):
    """添加带项目的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_page_number(slide, num, total):
    """添加页码"""
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f'{num} / {total}', font_size=10, color=GRAY,
                 alignment=PP_ALIGN.RIGHT)


# 计算总页数
TOTAL_SLIDES = 14

# ===========================
# SLIDE 1: 封面
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# 左半部分蓝色块
add_rect(slide, 0, 0, Inches(7), H, BLUE)

# 主标题
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(1.2),
             '医院志愿者助手', font_size=48, color=WHITE, bold=True)

# 副标题
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(0.8),
             '微信小程序 · 项目汇报', font_size=28, color=WHITE, bold=False)

# 分隔线
add_rect(slide, Inches(0.8), Inches(3.9), Inches(2), Inches(0.06), WHITE)

# 描述
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(1.2),
             '为医院志愿者打造的智能服务助手\n提升志愿服务效率与质量',
             font_size=18, color=RGBColor(0xBF, 0xDB, 0xFE))

# 右侧装饰 - 大图标文字
add_text_box(slide, Inches(8.5), Inches(2.5), Inches(4), Inches(3),
             '🏥', font_size=120, color=WHITE, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.5),
             '汇报日期：2025年 | 版本 v1.0', font_size=14, color=RGBColor(0x93, 0xC5, 0xFD))

add_page_number(slide, 1, TOTAL_SLIDES)


# ===========================
# SLIDE 2: 目录
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.8),
             '📋 汇报目录', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

toc_items = [
    '01    项目背景与目标',
    '02    核心功能展示',
    '03    技术架构与实现',
    '04    项目结构与代码组织',
    '05    AI 辅助开发经验',
    '06    未来迭代计划',
]

y_start = 2.2
for i, item in enumerate(toc_items):
    # 编号色块
    add_rect(slide, Inches(1.2), Inches(y_start + i * 0.85), Inches(0.06),
             Inches(0.5), BLUE)
    add_text_box(slide, Inches(1.6), Inches(y_start + i * 0.85), Inches(9), Inches(0.6),
                 item, font_size=22, color=DARK)

add_page_number(slide, 2, TOTAL_SLIDES)


# ===========================
# SLIDE 3: 项目背景
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.8),
             '🏥 项目背景', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 左侧：背景描述
bg_text = (
    '在医院环境中，志愿者经常需要解答患者的各种问题：\n\n'
    '• 某个症状应该挂哪个科室？\n'
    '• 各科室分布在医院的什么位置？\n'
    '• 医院的服务流程是怎样的？\n'
    '• 患者常问的问题如何快速回答？\n\n'
    '传统的纸质指南查阅不便，信息更新困难。\n'
    '一款便捷的移动端工具能大幅提升志愿者服务效率。'
)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(5),
             bg_text, font_size=16, color=DARK)

# 右侧：数据卡片
add_rect(slide, Inches(8.5), Inches(2.0), Inches(4), Inches(1.5), LIGHT_BLUE)
add_text_box(slide, Inches(8.8), Inches(2.3), Inches(3.5), Inches(0.8),
             '🎯 项目目标', font_size=20, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(8.8), Inches(3.0), Inches(3.5), Inches(0.5),
             '为志愿者提供一站式信息查询工具', font_size=14, color=GRAY)

add_rect(slide, Inches(8.5), Inches(3.8), Inches(4), Inches(1.5), RGBColor(0xD1, 0xFA, 0xE5))
add_text_box(slide, Inches(8.8), Inches(4.1), Inches(3.5), Inches(0.8),
             '👥 目标用户', font_size=20, color=GREEN, bold=True)
add_text_box(slide, Inches(8.8), Inches(4.8), Inches(3.5), Inches(0.5),
             '医院志愿者、导诊人员', font_size=14, color=GRAY)

add_rect(slide, Inches(8.5), Inches(5.6), Inches(4), Inches(1.5), RGBColor(0xFE, 0xF3, 0xC7))
add_text_box(slide, Inches(8.8), Inches(5.9), Inches(3.5), Inches(0.8),
             '⚡ 核心价值', font_size=20, color=ORANGE, bold=True)
add_text_box(slide, Inches(8.8), Inches(6.4), Inches(3.5), Inches(0.5),
             '快速响应患者问题，提升服务专业度', font_size=14, color=GRAY)

add_page_number(slide, 3, TOTAL_SLIDES)


# ===========================
# SLIDE 4: 功能总览
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.8),
             '🚀 核心功能总览', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

features = [
    ('🔍', '科室查询', '输入症状或疾病名称\n智能推荐对应就诊科室\n保存最近查询记录'),
    ('🗺️', '医院地图', '可视化楼层平面图\n各科室位置一目了然\n点击查看详细信息'),
    ('📚', '志愿者指南', '服务规范与注意事项\n应急处理流程指导\n沟通技巧与建议'),
    ('❓', '常见问题', '按分类浏览问题\n关键词快速搜索\n展开查看详细答案'),
    ('👤', '个人中心', '志愿者信息管理\n服务时长与次数统计\n服务记录与证书'),
]

for i, (icon, title, desc) in enumerate(features):
    x = Inches(0.6 + i * 2.5)
    # 卡片背景
    card = add_rect(slide, x, Inches(1.9), Inches(2.2), Inches(4.5), LIGHT_GRAY)
    card.shadow.inherit = False

    add_text_box(slide, x, Inches(2.1), Inches(2.2), Inches(0.8),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.9), Inches(2.2), Inches(0.5),
                 title, font_size=20, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7 + i * 2.5), Inches(3.5), Inches(2), Inches(2.5),
                 desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 4, TOTAL_SLIDES)


# ===========================
# SLIDE 5: 功能详情 - 科室查询
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '🔍 功能详情 — 科室查询', font_size=32, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 描述
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(12), Inches(0.6),
             '根据患者描述的症状或疾病名称，快速查询对应的就诊科室',
             font_size=16, color=GRAY)

# 流程步骤
steps = [
    ('1', '输入关键词', '患者描述症状\n如"头痛""发烧"', LIGHT_BLUE, BLUE),
    ('2', '系统匹配', '查询症状-科室\n映射数据库', RGBColor(0xD1, 0xFA, 0xE5), GREEN),
    ('3', '显示结果', '推荐就诊科室\n及科室介绍', RGBColor(0xFE, 0xF3, 0xC7), ORANGE),
    ('4', '历史记录', '保存查询记录\n方便快速复用', RGBColor(0xED, 0xE9, 0xFE), RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (num, title, desc, bg_c, border_c) in enumerate(steps):
    x = Inches(1.0 + i * 3.0)
    # 步骤卡片
    card = add_rect(slide, x, Inches(2.8), Inches(2.6), Inches(3.5), bg_c)
    # 步骤编号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(2.9), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = border_c
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    add_text_box(slide, x + Inches(0.1), Inches(3.8), Inches(2.4), Inches(0.5),
                 title, font_size=18, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(4.3), Inches(2.4), Inches(1.5),
                 desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# 技术要点
add_rect(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5), LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11), Inches(0.4),
             '💡 技术实现：症状-科室映射字典 + 本地存储缓存历史查询记录，无需后端即可运行',
             font_size=13, color=DARK_BLUE)

add_page_number(slide, 5, TOTAL_SLIDES)


# ===========================
# SLIDE 6: 功能详情 - 医院地图
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '🗺️ 功能详情 — 医院地图', font_size=32, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(12), Inches(0.6),
             '可视化楼层平面图，帮助志愿者快速定位各科室和设施位置',
             font_size=16, color=GRAY)

# 地图特性
features_data = [
    ('🏗️', '楼层切换', '支持 1-4 楼\n一键切换'),
    ('📐', '平面布局', '房间百分比定位\n模拟真实布局'),
    ('🎨', '颜色标识', '5 种颜色区分\n房间类型'),
    ('👆', '点击详情', '点击房间查看\n名称/介绍/位置'),
    ('📖', '图例说明', '底部图例辅助\n快速识别'),
]

for i, (icon, title, desc) in enumerate(features_data):
    x = Inches(0.6 + i * 2.5)
    add_rect(slide, x, Inches(2.6), Inches(2.2), Inches(3.2), LIGHT_GRAY)
    add_text_box(slide, x, Inches(2.8), Inches(2.2), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.5), Inches(2.2), Inches(0.5),
                 title, font_size=18, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7 + i * 2.5), Inches(4.0), Inches(2), Inches(1.5),
                 desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# 颜色图例
add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), LIGHT_GRAY)
colors = [('🟡', '服务窗口', '#F59E0B'), ('🔵', '门诊科室', '#3B82F6'),
          ('🟢', '专科门诊', '#10B981'), ('🔴', '住院/手术', '#EF4444'), ('🟣', '公共设施', '#8B5CF6')]
legend_text = '  '.join([f'{c[0]} {c[1]}' for c in colors])
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
             f'🎨 房间类型配色：  {legend_text}', font_size=13, color=GRAY)

add_page_number(slide, 6, TOTAL_SLIDES)


# ===========================
# SLIDE 7: 功能详情 - 志愿者指南 & 常见问题
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '📚 功能详情 — 志愿者指南 & 常见问题', font_size=32, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 左列 - 志愿者指南
add_rect(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(2.2), Inches(5), Inches(0.5),
             '📚 志愿者指南', font_size=22, color=DARK, bold=True)
add_rect(slide, Inches(1.0), Inches(2.7), Inches(1.2), Inches(0.04), BLUE)

guide_items = [
    '✅ 服务规范 — 着装要求、服务态度、服务内容',
    '✅ 应急处理 — 突发疾病、意外受伤、纠纷处理',
    '✅ 沟通技巧 — 倾听技巧、表达技巧、特殊患者沟通',
    '',
    '三 Tab 切换设计，图文并茂的指导内容',
    '帮助新志愿者快速上手',
]
add_bullet_list(slide, Inches(1.0), Inches(3.0), Inches(5), Inches(3.5),
                guide_items, font_size=14, color=DARK, spacing=Pt(6))

# 右列 - 常见问题
add_rect(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(0.5),
             '❓ 常见问题', font_size=22, color=DARK, bold=True)
add_rect(slide, Inches(7.0), Inches(2.7), Inches(1.2), Inches(0.04), BLUE)

faq_items = [
    '✅ 分类浏览 — 挂号缴费 / 就诊流程 / 医院设施',
    '✅ 关键词搜索 — 快速定位相关问题',
    '✅ 展开/收起 — 点击问题查看详细答案',
    '',
    '内置 8 条常见问题，覆盖患者高频咨询',
    '支持按分类和关键词双重过滤',
]
add_bullet_list(slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.5),
                faq_items, font_size=14, color=DARK, spacing=Pt(6))

add_page_number(slide, 7, TOTAL_SLIDES)


# ===========================
# SLIDE 8: 功能详情 - 个人中心
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '👤 功能详情 — 个人中心', font_size=32, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 三列卡片
cards = [
    ('📊', '服务统计', '服务时长：48 小时\n服务次数：126 次\n服务评分：5.0 🌟',
     RGBColor(0xDB, 0xEA, 0xFE), BLUE),
    ('📋', '功能菜单', '服务记录\n培训证书\n设置\n关于我们',
     RGBColor(0xD1, 0xFA, 0xE5), GREEN),
    ('🔐', '账户管理', '志愿者信息展示\n志愿者 ID 标识\n退出登录功能',
     RGBColor(0xFE, 0xF3, 0xC7), ORANGE),
]

for i, (icon, title, desc, bg, border) in enumerate(cards):
    x = Inches(1.0 + i * 4.0)
    add_rect(slide, x, Inches(2.0), Inches(3.5), Inches(4.5), bg)
    add_text_box(slide, x, Inches(2.2), Inches(3.5), Inches(0.7),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.0), Inches(3.5), Inches(0.5),
                 title, font_size=22, color=border, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.7), Inches(2.9), Inches(2.5),
                 desc, font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4), LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.82), Inches(11), Inches(0.35),
             '💡 注意：当前数据为静态示例，后续可接入云开发实现真实用户数据管理',
             font_size=12, color=DARK_BLUE)

add_page_number(slide, 8, TOTAL_SLIDES)


# ===========================
# SLIDE 9: 技术架构
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '🛠️ 技术架构', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 技术栈卡片
techs = [
    ('📱 前端框架', '微信小程序原生框架\nWXML + WXSS + JavaScript', LIGHT_BLUE, BLUE),
    ('🗄️ 数据存储', 'app.js 全局数据\nwx.setStorageSync 本地缓存', RGBColor(0xD1, 0xFA, 0xE5), GREEN),
    ('🎨 UI 设计', '自定义组件\nFlex 弹性布局\nrpx 响应式单位', RGBColor(0xFE, 0xF3, 0xC7), ORANGE),
    ('🛠️ 开发工具', '微信开发者工具\nTrae AI IDE\nVisual Studio Code', RGBColor(0xED, 0xE9, 0xFE), RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, desc, bg, border) in enumerate(techs):
    x = Inches(0.8 + i * 3.1)
    add_rect(slide, x, Inches(2.0), Inches(2.8), Inches(3.8), bg)
    add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(2.4), Inches(0.5),
                 title, font_size=18, color=border, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.4), Inches(2.5),
                 desc, font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)

# 架构特点
add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
             '✨ 架构特点：纯前端实现 · 零外部依赖 · 数据驱动视图 · 模块化页面设计 · 开箱即用',
             font_size=14, color=DARK, bold=False)

add_page_number(slide, 9, TOTAL_SLIDES)


# ===========================
# SLIDE 10: 项目结构
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '📁 项目结构', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 目录树 - 用代码块样式展示
tree_text = (
    'vollunteercloundapp/\n'
    '├── app.js           ← 全局入口 + 数据\n'
    '├── app.json         ← 全局配置 (4 Tab + 6 页面)\n'
    '├── app.wxss         ← 全局样式\n'
    '├── sitemap.json     ← 搜索配置\n'
    '├── project.config.json\n'
    '├── pages/\n'
    '│   ├── index/        ← 首页 (快捷入口)\n'
    '│   ├── deptSearch/   ← 科室查询\n'
    '│   ├── hospitalMap/  ← 医院地图 🆕\n'
    '│   ├── volunteerGuide/ ← 志愿者指南\n'
    '│   ├── faq/          ← 常见问题\n'
    '│   └── profile/      ← 个人中心\n'
    '├── 需求分析文档.md\n'
    '├── README.md\n'
    '└── generate_ppt.py  ← 本PPT生成脚本'
)

add_rect(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(5), RGBColor(0x1F, 0x29, 0x37))
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(6.6), Inches(4.8),
             tree_text, font_size=14, color=RGBColor(0xE5, 0xE7, 0xEB),
             font_name='Consolas')

# 右侧统计
stats_data = [
    ('6 个页面', '覆盖全部核心功能'),
    ('22 个房间', '4 层楼详细布局数据'),
    ('12+ 症状映射', '症状→科室智能推荐'),
    ('8 条常见问题', '覆盖高频咨询场景'),
    ('0 外部依赖', '纯原生实现'),
]

add_rect(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(5), LIGHT_GRAY)
add_text_box(slide, Inches(8.5), Inches(2.0), Inches(4), Inches(0.5),
             '📊 项目数据', font_size=20, color=DARK, bold=True)

for i, (num, label) in enumerate(stats_data):
    add_text_box(slide, Inches(8.5), Inches(2.7 + i * 0.85), Inches(4), Inches(0.4),
                 num, font_size=20, color=BLUE, bold=True)
    add_text_box(slide, Inches(8.5), Inches(3.1 + i * 0.85), Inches(4), Inches(0.4),
                 label, font_size=12, color=GRAY)

add_page_number(slide, 10, TOTAL_SLIDES)


# ===========================
# SLIDE 11: AI 辅助开发
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '🤖 AI 辅助开发经验', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 左侧：亮点
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
             '✅ AI 带来的优势', font_size=22, color=GREEN, bold=True)

pros = [
    '快速生成代码框架和页面结构',
    '完整的 UI 设计和交互逻辑',
    '数据结构和文档一次性生成',
    '开发效率大幅提升（估算节省 60%+ 时间）',
]
add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.5),
                pros, font_size=15, color=DARK, spacing=Pt(10))

# 右侧：问题与修复
add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.5),
             '⚠️ AI 生成需要注意的问题', font_size=22, color=ORANGE, bold=True)

cons = [
    '缺少 app = getApp() 引用 → 手动补全',
    '引用了不存在的页面路径 → 修正跳转',
    '图片资源路径缺失 → 补充资源文件',
    '空 navigator 链接 → 补充目标页面',
    '硬编码数据 → 需替换为动态数据源',
]
add_bullet_list(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(2.5),
                cons, font_size=15, color=DARK, spacing=Pt(8))

# 底部总结
add_rect(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), LIGHT_BLUE)
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(0.5),
             '💡 经验总结', font_size=20, color=DARK_BLUE, bold=True)
summary = (
    'AI 能快速产出 80% 的代码框架，但剩下的 20% 需要人工审查修正。\n'
    '关键在于：明确的需求描述 + 仔细的代码审查 + 针对性的手动优化。\n'
    'AI + 人工结合的开发模式，兼顾了效率与质量。'
)
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11), Inches(0.8),
             summary, font_size=14, color=DARK)

add_page_number(slide, 11, TOTAL_SLIDES)


# ===========================
# SLIDE 12: 后续计划
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '🔮 未来迭代计划', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

plans = [
    ('P0', '修复已知 Bug', '修复首页 app 引用、空跳转、缺失资源等问题', BLUE, '高优先'),
    ('P1', '接入腾讯地图', '使用真实地图 API 替代模拟楼层图', RGBColor(0x3B, 0x82, 0xF6), '中优先'),
    ('P1', '云开发接入', '微信云开发实现用户数据持久化存储', RGBColor(0x10, 0xB9, 0x81), '中优先'),
    ('P2', '扩充数据量', '丰富症状→科室库到 50+ 条映射', ORANGE, '低优先'),
    ('P2', '智能搜索', '支持模糊搜索和联想推荐', ORANGE, '低优先'),
    ('P2', '志愿者排班', '排班签到、服务记录统计功能', ORANGE, '低优先'),
]

for i, (priority, title, desc, color, level) in enumerate(plans):
    y = Inches(2.0 + i * 0.85)
    # 优先级标签
    tag = add_rect(slide, Inches(0.8), y, Inches(0.7), Inches(0.5), color)
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = priority
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    add_text_box(slide, Inches(1.7), y, Inches(3), Inches(0.5),
                 title, font_size=18, color=DARK, bold=True)
    add_text_box(slide, Inches(4.5), y, Inches(6), Inches(0.5),
                 desc, font_size=14, color=GRAY)
    # 优先级标签
    add_text_box(slide, Inches(11.5), y, Inches(1.2), Inches(0.5),
                 level, font_size=12, color=color, alignment=PP_ALIGN.RIGHT)

add_page_number(slide, 12, TOTAL_SLIDES)


# ===========================
# SLIDE 13: 数据看板
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
             '📊 项目数据看板', font_size=36, color=DARK, bold=True)
add_rect(slide, Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.05), BLUE)

# 数据卡片 - 4 个
metrics = [
    ('6', '页面总数', '覆盖全部核心功能模块', BLUE),
    ('4', 'Tab 导航', '首页 · 科室查询 · 地图 · 个人', GREEN),
    ('22', '地图房间', '4 层楼 · 5 种类型', ORANGE),
    ('1.5K+', '代码行数', 'JS + WXML + WXSS', RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (value, label, desc, color) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    # 大数字
    add_rect(slide, x, Inches(2.0), Inches(2.8), Inches(2.5), LIGHT_GRAY)
    add_text_box(slide, x, Inches(2.2), Inches(2.8), Inches(0.9),
                 value, font_size=48, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.2), Inches(2.8), Inches(0.4),
                 label, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.6), Inches(2.8), Inches(0.5),
                 desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# 底部时间线
add_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(0.4),
             '⏱️ 项目时间线', font_size=18, color=DARK, bold=True)

timeline = [
    ('需求分析', '明确功能定位\n与核心需求'),
    ('UI 设计', '界面原型\n与交互设计'),
    ('AI 编码', 'AI 辅助生成\n基础代码框架'),
    ('人工优化', 'Bug 修复\n地图可视化改造'),
    ('测试迭代', '功能测试\n体验优化'),
]

for i, (phase, desc) in enumerate(timeline):
    x = Inches(1.2 + i * 2.3)
    add_rect(slide, x, Inches(5.8), Inches(0.04), Inches(0.8), BLUE)
    # 圆点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.08), Inches(5.7), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE
    dot.line.fill.background()

    add_text_box(slide, x - Inches(0.5), Inches(6.0), Inches(2), Inches(0.4),
                 phase, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x - Inches(0.5), Inches(6.4), Inches(2), Inches(0.5),
                 desc, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 13, TOTAL_SLIDES)


# ===========================
# SLIDE 14: 致谢
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_rect(slide, 0, 0, W, H, BLUE)

add_text_box(slide, Inches(0), Inches(2.0), W, Inches(1.5),
             '感谢聆听', font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06), WHITE)

add_text_box(slide, Inches(0), Inches(3.9), W, Inches(0.8),
             '医院志愿者助手 · 微信小程序', font_size=24, color=RGBColor(0xBF, 0xDB, 0xFE),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.0), W, Inches(0.8),
             '有任何问题欢迎交流探讨 🙋', font_size=18, color=RGBColor(0x93, 0xC5, 0xFD),
             alignment=PP_ALIGN.CENTER)

add_page_number(slide, 14, TOTAL_SLIDES)


# ========== 保存 ==========
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           '医院志愿者助手_项目汇报.pptx')
prs.save(output_path)
print(f'✅ PPT 已生成: {output_path}')
print(f'📐 共 {len(prs.slides)} 页幻灯片')
